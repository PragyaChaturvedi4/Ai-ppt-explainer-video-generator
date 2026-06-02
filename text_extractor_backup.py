from pptx import Presentation
import fitz
import os

def extract_text_from_ppt(path):
    if not os.path.exists(path):
        raise FileNotFoundError(f"PPT not found: {path}")

    prs = Presentation(path)
    slides_text = []

    for slide in prs.slides:
        content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                content.append(shape.text)

        slides_text.append("\n".join(content))

    return slides_text


def extract_text_from_pdf(path):
    if not os.path.exists(path):
        raise FileNotFoundError(f"PDF not found: {path}")

    doc = fitz.open(path)
    pages = []

    for page in doc:
        pages.append(page.get_text())

    doc.close()
    return pages

